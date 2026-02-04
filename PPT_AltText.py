import sys, os, argparse
from pathlib import Path
from dotenv import load_dotenv
from pptx import Presentation
from google import genai

# Load API Key
load_dotenv()
client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))

def get_alt_text(image_bytes, discipline):
    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[
                f"Act as an expert in {discipline}. Provide a concise, 1-sentence accessibility alt-text.",
                genai.types.Part.from_bytes(data=image_bytes, mime_type='image/png')
            ]
        )
        # Clean disclaimers
        text = response.text.strip()
        for tag in ["Text automatically generated.", "Description automatically generated."]:
            text = text.replace(tag, "")
        return text.strip()
    except: return "Biochemical diagram."

def apply_text_robustly(shape, text):
    """Writes to every metadata field to ensure UI visibility."""
    # Method 1: Standard library properties
    try:
        shape.name = text
        shape.description = text
    except: pass

    # Method 2: Force XML 'descr' (The UI Pane's primary target)
    try:
        # Check both common namespaces for pictures and generic shapes
        namespaces = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }
        
        # Navigate to the non-visual properties element
        nv_props = shape._element.find('.//p:cNvPr', namespaces)
        if nv_props is not None:
            nv_props.set('descr', text)
        else:
            # Try the alternative picture-specific path
            nv_pic_props = shape._element.find('.//p:nvPicPr/p:cNvPicPr', namespaces)
            if nv_pic_props is not None:
                nv_pic_props.set('descr', text)
    except: pass

def force_tag(shape, discipline):
    count = 0
    # Process images, placeholders with images, and nested groups
    if hasattr(shape, "image"):
        alt = get_alt_text(shape.image.blob, discipline)
        apply_text_robustly(shape, alt)
        count += 1
    elif shape.shape_type == 6: # Grouped Object
        for sub_shape in shape.shapes:
            count += force_tag(sub_shape, discipline)
    return count

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("target", help="File or folder to process")
    parser.add_argument("-d", "--discipline", default="Biochemistry")
    args = parser.parse_args()
    
    target = Path(args.target)
    files = [target] if target.is_file() else list(target.glob("*.pptx"))
    
    for f in files:
        if "_Accessible" in f.name or f.name.startswith((".","~$")): continue
        print(f"Deep-Scrubbing: {f.name}...")
        try:
            prs = Presentation(f)
            total = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    total += force_tag(shape, args.discipline)
            
            out = f.parent / f"{f.stem}_Accessible{f.suffix}"
            prs.save(out)
            print(f"Success: {total} items tagged in {out.name}\n")
        except Exception as e:
            print(f"Error processing {f.name}: {e}")

if __name__ == "__main__": main()